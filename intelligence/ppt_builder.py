from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(11, 16, 32)
HEADER = RGBColor(7, 11, 22)
PANEL = RGBColor(17, 24, 39)
PANEL_2 = RGBColor(23, 32, 51)
TEXT = RGBColor(229, 231, 235)
MUTED = RGBColor(156, 163, 175)
BLUE = RGBColor(56, 189, 248)
GOLD = RGBColor(214, 168, 92)
GREEN = RGBColor(52, 211, 153)
RED = RGBColor(248, 113, 113)
LINE = RGBColor(51, 65, 85)
WHITE = RGBColor(255, 255, 255)


def _items(value: Any) -> list[Any]:
    return value if isinstance(value, list) else []


def _dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _text(value: Any, fallback: str = "今日无可靠更新") -> str:
    text = str(value or "").strip()
    return text or fallback


def _short(value: Any, limit: int = 74) -> str:
    text = _text(value, "")
    return text if len(text) <= limit else text[: limit - 1] + "…"


def _run(paragraph, text: str, size: int, color: RGBColor = TEXT, bold: bool = False) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei UI"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def _box(slide, left, top, width, height, text: str, size: int = 12, color: RGBColor = TEXT, bold: bool = False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    p = frame.paragraphs[0]
    p.alignment = align
    _run(p, text, size, color, bold)
    return shape


def _fill_rect(slide, left, top, width, height, color: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.6)
    else:
        shape.line.fill.background()
    return shape


def _blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def _title(slide, page: int, kicker: str, title: str) -> None:
    _fill_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.72), HEADER)
    _box(slide, Inches(0.42), Inches(0.18), Inches(2.2), Inches(0.22), kicker, 8, GOLD, True)
    _box(slide, Inches(0.42), Inches(0.40), Inches(10.6), Inches(0.36), title, 18, TEXT, True)
    _fill_rect(slide, Inches(0.42), Inches(0.77), Inches(0.9), Inches(0.035), GOLD)
    _box(slide, Inches(11.95), Inches(6.98), Inches(0.7), Inches(0.2), f"{page:02d}", 8, MUTED, False, PP_ALIGN.RIGHT)
    _box(slide, Inches(0.42), Inches(6.96), Inches(8.8), Inches(0.22), "真实新闻源驱动 | 无可靠来源不纳入正文 | 不构成投资建议", 7, MUTED)


def _panel(slide, left, top, width, height, title: str, body: str, accent: RGBColor = BLUE) -> None:
    _fill_rect(slide, left, top, width, height, PANEL, LINE)
    _fill_rect(slide, left, top, Inches(0.04), height, accent)
    _box(slide, left + Inches(0.12), top + Inches(0.09), width - Inches(0.24), Inches(0.24), title, 9, accent, True)
    _box(slide, left + Inches(0.12), top + Inches(0.39), width - Inches(0.24), height - Inches(0.47), body, 8, TEXT)


def _bullets(slide, left, top, width, height, items: list[Any], max_items: int = 5, size: int = 12) -> None:
    selected = items[:max_items] or ["今日无可靠更新"]
    text = "\n".join(f"• {_short(item, 96)}" for item in selected)
    _box(slide, left, top, width, height, text, size, TEXT)


def _table(slide, left, top, width, height, headers: list[str], rows: list[list[str]], font_size: int = 6) -> None:
    if not rows:
        rows = [["今日无可靠更新"] + [""] * (len(headers) - 1)]
    row_count = len(rows[:10]) + 1
    table = slide.shapes.add_table(row_count, len(headers), left, top, width, height).table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL_2
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Microsoft YaHei UI"
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = BLUE
    for row_index, row in enumerate(rows[:10], start=1):
        for col, value in enumerate(row[: len(headers)]):
            cell = table.cell(row_index, col)
            cell.text = _short(value, 60)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if row_index % 2 else BG
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Microsoft YaHei UI"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = TEXT


def _cover(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _fill_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(6, 10, 20))
    for idx, color in enumerate([BLUE, GOLD, GREEN]):
        _fill_rect(slide, Inches(8.6 + idx * 1.05), Inches(1.0 + idx * 0.55), Inches(0.55), Inches(4.25 - idx * 0.55), color)
    _box(slide, Inches(0.65), Inches(1.25), Inches(7.5), Inches(0.7), "全球科技资本市场日报", 30, TEXT, True)
    _box(slide, Inches(0.7), Inches(2.05), Inches(7.4), Inches(0.35), "Technology & Capital Market Daily Briefing", 15, BLUE)
    _box(slide, Inches(0.7), Inches(2.75), Inches(3.0), Inches(0.35), _text(report.get("date")), 14, GOLD, True)
    _box(slide, Inches(0.7), Inches(6.15), Inches(9.2), Inches(0.45), "仅收录过去24小时内可追溯来源；未验证信息不纳入正文。", 11, MUTED)


def _core(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 2, "EXECUTIVE SUMMARY", "今日核心结论：只输出可验证事件驱动的资本市场判断")
    _bullets(slide, Inches(0.62), Inches(1.25), Inches(6.9), Inches(4.2), _items(report.get("core_conclusions")), 5, 13)
    flows = _items(report.get("capital_flows"))
    for idx, item in enumerate(flows[:3]):
        item = _dict(item)
        _panel(slide, Inches(8.0), Inches(1.25 + idx * 1.28), Inches(3.45), Inches(0.85), _text(item.get("sector")), _short(item.get("outlook"), 90), [BLUE, RED, GOLD][idx])


def _top_news(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 3, "MARKET MOVING EVENTS", "全球科技新闻 Top10：按资本市场影响排序")
    rows = []
    for item in _items(report.get("top_news"))[:10]:
        item = _dict(item)
        rows.append([str(item.get("rank", "")), _text(item.get("title")), _text(item.get("source")), _text(item.get("sector")), _text(item.get("sentiment")), _text(item.get("bank_view"))])
    _table(slide, Inches(0.45), Inches(1.15), Inches(12.4), Inches(5.35), ["排名", "事件", "来源", "板块", "方向", "投行视角"], rows, 5)


def _capital_flows(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 4, "CAPITAL FLOWS", "全球资本流向：赛道热度、资金关注与驱动事件")
    flows = _items(report.get("capital_flows"))
    for idx, item in enumerate(flows[:8]):
        item = _dict(item)
        left = Inches(0.55 + (idx % 2) * 6.25)
        top = Inches(1.15 + (idx // 2) * 1.28)
        _panel(slide, left, top, Inches(5.35), Inches(0.9), f"{_text(item.get('sector'))} | {item.get('heat', '')}", _short(f"{item.get('driver', '')}；{item.get('outlook', '')}", 120), GOLD if idx < 2 else BLUE)


def _deep_dive(prs: Presentation, report: dict[str, Any], index: int, page: int) -> None:
    item = _dict((_items(report.get("deep_dives")) + [{}] * 3)[index])
    slide = _blank(prs)
    _title(slide, page, "DEEP DIVE", f"深度分析 {index + 1}：{_short(item.get('event'), 56)}")
    panels = [
        ("发生了什么", item.get("what_happened"), GOLD),
        ("为什么重要", item.get("why_important"), BLUE),
        ("产业链影响", item.get("supply_chain_impact"), GREEN),
        ("收入 / 利润 / 估值", f"{item.get('revenue_impact', '')} {item.get('profit_impact', '')} {item.get('valuation_impact', '')}", BLUE),
        ("交易观察", item.get("short_term_opportunity"), GOLD),
        ("核心风险", item.get("core_risk"), RED),
    ]
    for idx, (title, body, accent) in enumerate(panels):
        _panel(slide, Inches(0.6 + (idx % 3) * 4.15), Inches(1.2 + (idx // 3) * 2.05), Inches(3.55), Inches(1.25), title, _text(body), accent)


def _leader_views(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 8, "LEADER COMMENTARY", "科技大佬观点追踪：公开出处、中文翻译与可信度评分")
    rows = []
    for item in _items(report.get("leader_views"))[:8]:
        item = _dict(item)
        rows.append([_text(item.get("person")), _text(item.get("latest_comment")), _text(item.get("source")), _text(item.get("attention")), _text(item.get("bank_view"))])
    _table(slide, Inches(0.55), Inches(1.15), Inches(12.0), Inches(5.35), ["人物", "最新发言", "出处", "关注度", "投行解读"], rows, 5)


def _mega_caps(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 9, "MEGA CAP TECH", "美股科技巨头追踪：新闻、情绪、短期和长期影响")
    rows = []
    for item in _items(report.get("mega_cap_tracking"))[:12]:
        item = _dict(item)
        rows.append([_text(item.get("ticker")), _text(item.get("news")), _text(item.get("sentiment")), _text(item.get("short_term")), _text(item.get("long_term")), _text(item.get("rating"))])
    _table(slide, Inches(0.45), Inches(1.15), Inches(12.4), Inches(5.35), ["Ticker", "今日新闻", "情绪", "短期", "长期", "等级"], rows, 5)


def _ai_chain(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 10, "AI VALUE CHAIN", "AI产业链跟踪：模型、算力、芯片、应用的受益链")
    rows = []
    for item in _items(report.get("ai_chain")):
        item = _dict(item)
        rows.append([_text(item.get("layer")), _text(item.get("beneficiaries")), _text(item.get("losers")), _text(item.get("trend"))])
    _table(slide, Inches(0.65), Inches(1.2), Inches(11.6), Inches(4.6), ["环节", "谁受益", "谁受损", "趋势"], rows, 6)


def _crypto(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 11, "CRYPTO MARKETS", "加密货币市场跟踪：监管、ETF资金流和机构入场")
    crypto = _dict(report.get("crypto"))
    rows = [[key, _text(value)] for key, value in crypto.items()]
    _table(slide, Inches(0.9), Inches(1.15), Inches(10.9), Inches(5.25), ["主题", "今日更新"], rows, 6)


def _opportunity_risk(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 12, "WATCHLIST / RISKS", "交易机会观察池与风险预警：只给观察逻辑")
    pool = _dict(report.get("opportunity_pool"))
    for idx, (key, label, color) in enumerate([("short_term", "短线投机机会", GOLD), ("mid_term", "中线波段机会", BLUE), ("long_term", "长线投资机会", GREEN)]):
        first = _dict((_items(pool.get(key)) + [{}])[0])
        body = f"标的：{first.get('target', '')}\n逻辑：{first.get('logic', '')}\n催化剂：{first.get('catalyst', '')}\n风险：{first.get('risk', '')}"
        _panel(slide, Inches(0.6 + idx * 4.15), Inches(1.2), Inches(3.55), Inches(2.1), label, body, color)
    _panel(slide, Inches(0.6), Inches(4.0), Inches(11.85), Inches(1.05), "风险预警", "；".join(str(item) for item in _items(report.get("risk_alerts"))[:8]), RED)


def _next_sources(prs: Presentation, report: dict[str, Any]) -> None:
    slide = _blank(prs)
    _title(slide, 13, "NEXT 72H / SOURCES", "未来72小时关注清单与信息来源汇总")
    next_72 = _dict(report.get("next_72h"))
    rows = []
    for key in ["24h", "48h", "72h"]:
        rows.append([key, "；".join(str(item) for item in _items(next_72.get(key))[:4])])
    _table(slide, Inches(0.65), Inches(1.15), Inches(11.6), Inches(1.8), ["时间", "关注事项"], rows, 7)
    source_rows = []
    for item in _items(report.get("sources"))[:8]:
        item = _dict(item)
        source_rows.append([_text(item.get("source")), _short(item.get("title"), 80), _text(item.get("published_at")), _text(item.get("url"), "")])
    _table(slide, Inches(0.65), Inches(3.55), Inches(11.6), Inches(2.35), ["来源", "标题", "时间", "链接"], source_rows, 5)


def build_pptx(report: dict[str, Any], output_path: str | Path) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _cover(prs, report)
    _core(prs, report)
    _top_news(prs, report)
    _capital_flows(prs, report)
    _deep_dive(prs, report, 0, 5)
    _deep_dive(prs, report, 1, 6)
    _deep_dive(prs, report, 2, 7)
    _leader_views(prs, report)
    _mega_caps(prs, report)
    _ai_chain(prs, report)
    _crypto(prs, report)
    _opportunity_risk(prs, report)
    _next_sources(prs, report)

    prs.save(output_path)
    return output_path
